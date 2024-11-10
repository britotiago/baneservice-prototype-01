# file_extractors.py
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
import PyPDF2
from PIL import Image
import os

# DOCX Extractor
def extract_text_from_docx(file_path):
    doc = Document(file_path)
    text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
    return text

# PPTX Extractor
def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return '\n'.join(text)

# XLSX Extractor
def extract_text_from_xlsx(file_path):
    workbook = load_workbook(filename=file_path)
    text = []
    for sheet in workbook:
        for row in sheet.iter_rows(values_only=True):
            text.append(' '.join([str(cell) for cell in row if cell is not None]))
    return '\n'.join(text)

# PDF Extractor
def extract_text_from_pdf(file_path):
    with open(file_path, 'rb') as pdf_file:
        reader = PyPDF2.PdfReader(pdf_file)
        text = ''
        for page in reader.pages:
            text += page.extract_text()
    return text

# Image (JPG/PNG) Extractor
def extract_text_from_image(file_path):
    img = Image.open(file_path)
    text = pytesseract.image_to_string(img)
    return text

# Main function to handle different file types
def extract_text_from_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == '.docx':
        return extract_text_from_docx(file_path)
    elif ext == '.pptx':
        return extract_text_from_pptx(file_path)
    elif ext == '.xlsx':
        return extract_text_from_xlsx(file_path)
    elif ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif ext in ['.jpg', '.png']:
        return extract_text_from_image(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

def chunk_text(text, max_tokens=7500):
    words = text.split()
    chunks = []
    current_chunk = []

    for word in words:
        if len(current_chunk) + len(word.split()) <= max_tokens:
            current_chunk.append(word)
        else:
            chunks.append(' '.join(current_chunk))
            current_chunk = [word]
    if current_chunk:
        chunks.append(' '.join(current_chunk))
    return chunks

def process_files_in_directory(directory):
    file_summaries = []
    for file_name in os.listdir(directory):
        file_path = os.path.join(directory, file_name)
        try:
            file_text = extract_text_from_file(file_path)
            chunks = chunk_text(file_text)
            file_summaries.append({
                "file_name": file_name,
                "chunks": chunks
            })
        except Exception as e:
            print(f"Failed to process file {file_name}: {e}")
    return file_summaries